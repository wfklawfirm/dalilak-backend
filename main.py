#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Dalilak AI — FastAPI Backend v4 (Auth + Subscriptions + Admin)"""

import asyncio
import base64
import hashlib
import io
import json
import logging
import os
import secrets
import sys
import time
import uuid
from collections import OrderedDict
from contextvars import ContextVar
from datetime import datetime, timedelta, timezone
from typing import AsyncGenerator, Optional

import httpx
import jwt as _jwt
from fastapi import Depends, FastAPI, Header, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from openai import AsyncOpenAI
from pydantic import BaseModel
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance, FieldCondition, Filter, MatchValue,
    PayloadSchemaType, PointIdsList, PointStruct, VectorParams,
)

# ═══════════════════════════════════════════════════════════════
#  DOCUMENT EXTRACTION HELPERS
# ═══════════════════════════════════════════════════════════════

def extract_text_from_pdf(b64: str) -> str:
    try:
        import pdfplumber
        raw = base64.b64decode(b64)
        parts = []
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page in pdf.pages[:20]:
                t = page.extract_text()
                if t:
                    parts.append(t)
        return "\n\n".join(parts)[:15000]
    except Exception:
        try:
            import fitz
            raw = base64.b64decode(b64)
            doc = fitz.open(stream=raw, filetype="pdf")
            parts = [doc[i].get_text() for i in range(min(20, len(doc)))]
            doc.close()
            return "\n\n".join(parts)[:15000]
        except Exception as e:
            return f"[تعذّر استخراج نص PDF: {e}]"

def extract_text_from_docx(b64: str) -> str:
    try:
        from docx import Document
        raw = base64.b64decode(b64)
        doc = Document(io.BytesIO(raw))
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                lines.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(lines)[:15000]
    except Exception as e:
        return f"[تعذّر استخراج نص Word: {e}]"

def extract_text_from_excel(b64: str) -> str:
    try:
        import openpyxl
        raw = base64.b64decode(b64)
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"--- ورقة: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in vals):
                    lines.append(" | ".join(vals))
        return "\n".join(lines)[:15000]
    except Exception as e:
        return f"[تعذّر استخراج Excel: {e}]"

def extract_text_from_pptx(b64: str) -> str:
    try:
        from pptx import Presentation
        raw = base64.b64decode(b64)
        prs = Presentation(io.BytesIO(raw))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"--- شريحة {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)[:15000]
    except Exception as e:
        return f"[تعذّر استخراج PowerPoint: {e}]"

def extract_text_from_csv(b64: str) -> str:
    try:
        import csv
        raw = base64.b64decode(b64).decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(raw))
        lines = [" | ".join(row) for row in reader if any(c.strip() for c in row)]
        return "\n".join(lines[:500])[:15000]
    except Exception as e:
        return f"[تعذّر قراءة CSV: {e}]"

def extract_text_from_zip(b64: str) -> str:
    try:
        import zipfile
        raw = base64.b64decode(b64)
        text_exts = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".py", ".js", ".ts"}
        lines = []
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            lines.append(f"محتويات الأرشيف ({len(zf.namelist())} ملف):")
            for name in zf.namelist()[:50]:
                lines.append(f"  - {name}")
            lines.append("")
            for name in zf.namelist()[:10]:
                ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""
                if ext in text_exts:
                    try:
                        content = zf.read(name).decode("utf-8", errors="replace")[:3000]
                        lines.append(f"=== {name} ===\n{content}")
                    except Exception:
                        pass
        return "\n".join(lines)[:15000]
    except Exception as e:
        return f"[تعذّر فتح ZIP: {e}]"

async def transcribe_audio(b64: str, file_type: str, file_name: str) -> str:
    try:
        raw = base64.b64decode(b64)
        ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else "mp3"
        fname = f"audio.{ext}"
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        try:
            with open(tmp_path, "rb") as f:
                transcript = await oai().audio.transcriptions.create(
                    model="whisper-1", file=(fname, f, file_type or "audio/mpeg"),
                    response_format="text"
                )
            return str(transcript)[:15000]
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return f"[تعذّر تحويل الصوت: {e}]"

# ═══════════════════════════════════════════════════════════════
#  CONFIG
# ═══════════════════════════════════════════════════════════════

COLLECTION     = "dalilak_ai_v2"
EMBED_MODEL    = "text-embedding-3-large"
VECTOR_DIM     = 3072
MODEL_FAST     = "gpt-4o-mini"
MODEL_SMART    = "gpt-4o"
MIN_SCORE           = 0.28   # Qdrant retrieval floor (chunks below this are never returned)
SUFFICIENCY_TOP_SCORE = 0.35  # Phase 5 gate: top chunk must reach this to proceed to GPT
MAX_CTX        = 12
MAX_TOKENS     = 3000
MAX_HISTORY    = 6
MAX_CHARS      = 12000
MAX_DOC_TOKENS = 3500

# Phase 12 — Performance hardening: timeout ceilings for external calls
QDRANT_TIMEOUT_SEC  = 10    # Qdrant vector search must complete within 10 s
OPENAI_TIMEOUT_SEC  = 60    # OpenAI completions must complete within 60 s
MAX_MESSAGE_LEN     = 4000  # Max chars in a single user message (prevents token flood)

# Auth config
JWT_SECRET   = os.environ.get("JWT_SECRET", "dalilak-secret-CHANGE-IN-PROD")
JWT_ALGO     = "HS256"
TRIAL_DAYS   = 3
ADMIN_SECRET      = os.environ.get("ADMIN_SECRET", "dalilak-admin-CHANGE-IN-PROD")
# ── Email / Password-reset ────────────────────────────────────────────────────
# SECURITY: RESEND_API_KEY is consumed inside email_service.py only.
# It is never read, logged, or surfaced here.
RESEND_FROM_EMAIL = os.environ.get("RESEND_FROM_EMAIL", "noreply@dalilak.ai")
APP_BASE_URL      = os.environ.get("APP_BASE_URL", "https://dalilak-frontend.vercel.app").rstrip("/")

from rate_limit    import enforce         as _rate_enforce     # noqa: E402
from email_service import send_reset_email as _send_reset_email  # noqa: E402
from plan_quota    import check_and_increment as _check_quota   # Phase 10  # noqa: E402

# ── JWT_SECRET startup validation ─────────────────────────────────────────────
# Must fire before any external service client (Qdrant, OpenAI, …) is created.
# SECURITY: JWT_SECRET value is never logged, printed, or included in output.
from config import validate_security_configuration as _validate_cfg, ConfigurationError as _CfgError  # noqa: E402
try:
    _validate_cfg(JWT_SECRET)
except _CfgError as _exc:
    sys.stderr.write(str(_exc) + "\n")
    sys.exit(1)
# ─────────────────────────────────────────────────────────────────────────────

# Qdrant collections for users & logs
USERS_COL    = "dalilak_users"
LOGS_COL     = "dalilak_logs"
RESETS_COL   = "dalilak_resets"
_users_ready = _logs_ready = _resets_ready = False

# ═══════════════════════════════════════════════════════════════
#  SYSTEM PROMPT
# ═══════════════════════════════════════════════════════════════

_PROMPT_PATH = os.path.join(os.path.dirname(__file__), "system_prompt.txt")
try:
    with open(_PROMPT_PATH, encoding="utf-8") as f:
        SYSTEM_PROMPT = f.read()
except Exception:
    SYSTEM_PROMPT = "أنت دليلك AI، مساعد المواطن اللبناني في كل الشؤون الحكومية."

# ═══════════════════════════════════════════════════════════════
#  LAZY CLIENTS
# ═══════════════════════════════════════════════════════════════

_oai: Optional[AsyncOpenAI] = None
_qdrant: Optional[QdrantClient] = None

def oai() -> AsyncOpenAI:
    global _oai
    if _oai is None:
        _oai = AsyncOpenAI(api_key=os.environ.get("OPENAI_API_KEY", ""))
    return _oai

def qdrant() -> QdrantClient:
    global _qdrant
    if _qdrant is None:
        _qdrant = QdrantClient(
            url=os.environ.get("QDRANT_URL", "").rstrip("/"),
            api_key=os.environ.get("QDRANT_API_KEY"),
            timeout=30,
        )
    return _qdrant

def qdrant_url() -> str:
    return os.environ.get("QDRANT_URL", "").rstrip("/")

def qdrant_headers() -> dict:
    return {"api-key": os.environ.get("QDRANT_API_KEY", ""), "Content-Type": "application/json"}

# ═══════════════════════════════════════════════════════════════
#  LRU ANSWER CACHE
# ═══════════════════════════════════════════════════════════════

_CACHE_MAX = 200
_cache: OrderedDict[str, dict] = OrderedDict()

def _ck(q: str, d: Optional[str]) -> str:
    return hashlib.md5(f"{q.strip().lower()}||{d or ''}".encode()).hexdigest()

def _cget(key: str) -> Optional[dict]:
    v = _cache.get(key)
    if v:
        _cache.move_to_end(key)
    return v

def _cset(key: str, val: dict) -> None:
    _cache[key] = val
    _cache.move_to_end(key)
    while len(_cache) > _CACHE_MAX:
        _cache.popitem(last=False)

# ═══════════════════════════════════════════════════════════════
#  PASSWORD HELPERS
# ═══════════════════════════════════════════════════════════════

def hash_pw(password: str) -> str:
    salt = secrets.token_hex(16)
    key = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 100_000)
    return f"{salt}:{key.hex()}"

def verify_pw(password: str, stored: str) -> bool:
    try:
        salt, key_hex = stored.split(":", 1)
        key = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 100_000)
        return key.hex() == key_hex
    except Exception:
        return False

def _hash_reset_token(token: str) -> str:
    """
    SHA-256 digest of a raw reset token.
    Only the hash is stored server-side; the raw token travels to the user's
    inbox only and is never persisted.
    """
    return hashlib.sha256(token.encode()).hexdigest()

# ═══════════════════════════════════════════════════════════════
#  PHASE 9 — OBSERVABILITY + PRIVACY
# ═══════════════════════════════════════════════════════════════

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [dalilak] %(message)s",
    datefmt="%Y-%m-%dT%H:%M:%SZ",
)
_log = logging.getLogger("dalilak")

# Per-request correlation ID — propagated to X-Request-ID response header.
# Never used to log user query content.
_req_id_var: ContextVar[str] = ContextVar("req_id", default="-")

# ═══════════════════════════════════════════════════════════════
#  JWT HELPERS
# ═══════════════════════════════════════════════════════════════

JWT_EXPIRY_DAYS = 7   # Phase 8: reduced from 30 days

# Phase 8 — In-memory token blocklist: {jti: exp_epoch_float}
# Cleared on restart; safe because JWTs also expire in 7 days.
_revoked_tokens: dict[str, float] = {}

def _blocklist_prune() -> None:
    """Evict expired entries — they can no longer be used anyway."""
    now = time.time()
    stale = [jti for jti, exp in _revoked_tokens.items() if exp < now]
    for jti in stale:
        _revoked_tokens.pop(jti, None)

# ── Redis JTI state machine ────────────────────────────────────────────────
# Three states, never conflated:
#   REDIS_NOT_CONFIGURED        — REDIS_URL absent at startup; in-memory only
#   REDIS_HEALTHY               — connected and responding
#   REDIS_TEMPORARILY_UNAVAILABLE — was healthy, now unreachable
_REDIS_NOT_CONFIGURED = "REDIS_NOT_CONFIGURED"
_REDIS_HEALTHY        = "REDIS_HEALTHY"
_REDIS_UNAVAILABLE    = "REDIS_TEMPORARILY_UNAVAILABLE"

_redis_jti_status: str = _REDIS_NOT_CONFIGURED   # updated at first use
_REDIS_JTI_TIMEOUT = 2.0  # seconds — short timeout for auth-path Redis calls

import os as _os
_REDIS_CONFIGURED_AT_STARTUP: bool = bool(_os.environ.get("REDIS_URL", "").strip())
if not _REDIS_CONFIGURED_AT_STARTUP:
    _log.warning(
        "[jti] REDIS_URL not set — JTI revocation is in-memory only. "
        "Logout revocation will NOT survive a backend restart. "
        "Set REDIS_URL to enable durable revocation."
    )

async def _jti_revoke_redis(jti: str, exp_ts: float) -> None:
    """
    Write revoked JTI to Redis with remaining-lifetime TTL.
    If Redis is not configured, silently skips (in-memory path is sufficient).
    If Redis is configured but fails, logs operational alert — does NOT silently succeed.
    Privacy: jti value is never logged.
    """
    global _redis_jti_status
    if not _REDIS_CONFIGURED_AT_STARTUP:
        return  # in-memory path; degraded health already logged at startup
    try:
        from rate_limit import _get_redis
        r = await _get_redis()
        if r is None:
            _redis_jti_status = _REDIS_UNAVAILABLE
            _log.error(
                "[jti][ALERT] Redis configured but client unavailable during logout revocation. "
                "Revocation stored in-memory only for this instance. "
                "req_id=<see request log>"
            )
            return
        ttl = max(1, int(exp_ts - time.time()))
        await asyncio.wait_for(
            r.set(f"dalilak:jti:{jti}", "1", ex=ttl),
            timeout=_REDIS_JTI_TIMEOUT,
        )
        _redis_jti_status = _REDIS_HEALTHY
    except asyncio.TimeoutError:
        _redis_jti_status = _REDIS_UNAVAILABLE
        _log.error(
            "[jti][ALERT] Redis timeout during logout revocation (%.1fs limit). "
            "Revocation in-memory only for this instance.",
            _REDIS_JTI_TIMEOUT,
        )
    except Exception as exc:
        _redis_jti_status = _REDIS_UNAVAILABLE
        _log.error("[jti][ALERT] Redis error during logout revocation: %s", type(exc).__name__)


async def _jti_is_revoked_redis(jti: str) -> bool:
    """
    Check Redis for JTI revocation.

    Policy:
    - REDIS_NOT_CONFIGURED  → return False (in-memory is the only store; already checked)
    - REDIS_HEALTHY/unknown → query Redis; raise 503 on failure (do NOT accept token)
    - timeout               → raise 503 (fail closed when Redis was previously healthy)

    Privacy: jti value is never logged.
    Raises HTTPException(503) if Redis is configured but unreachable.
    """
    global _redis_jti_status
    if not _REDIS_CONFIGURED_AT_STARTUP:
        return False  # in-memory already checked by caller

    try:
        from rate_limit import _get_redis
        r = await _get_redis()
        if r is None:
            # Redis was configured but client unavailable — fail closed
            _redis_jti_status = _REDIS_UNAVAILABLE
            _log.error(
                "[jti][ALERT] Redis configured but client unavailable during token verification."
            )
            raise HTTPException(
                503,
                detail="خدمة التحقق من الجلسة غير متاحة مؤقتاً — حاول مجدداً",
            )
        result = await asyncio.wait_for(
            r.exists(f"dalilak:jti:{jti}"),
            timeout=_REDIS_JTI_TIMEOUT,
        )
        _redis_jti_status = _REDIS_HEALTHY
        return result > 0
    except HTTPException:
        raise
    except asyncio.TimeoutError:
        _redis_jti_status = _REDIS_UNAVAILABLE
        _log.error(
            "[jti][ALERT] Redis timeout during token verification (%.1fs limit).",
            _REDIS_JTI_TIMEOUT,
        )
        raise HTTPException(
            503,
            detail="خدمة التحقق من الجلسة غير متاحة مؤقتاً — حاول مجدداً",
        )
    except Exception as exc:
        _redis_jti_status = _REDIS_UNAVAILABLE
        _log.error("[jti][ALERT] Redis error during token verification: %s", type(exc).__name__)
        raise HTTPException(
            503,
            detail="خدمة التحقق من الجلسة غير متاحة مؤقتاً — حاول مجدداً",
        )

def create_token(username: str, role: str = "user") -> str:
    now = datetime.now(timezone.utc)
    exp = now + timedelta(days=JWT_EXPIRY_DAYS)
    payload = {
        "iss": "dalilak-ai",
        "sub": username,
        "role": role,
        "jti": secrets.token_hex(8),   # unique token ID — used for blocklist
        "iat": int(now.timestamp()),
        "exp": int(exp.timestamp()),
    }
    return _jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGO)

def decode_token(token: str) -> dict:
    return _jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGO])

# ═══════════════════════════════════════════════════════════════
#  QDRANT USER STORE
# ═══════════════════════════════════════════════════════════════

def _ensure_users() -> None:
    global _users_ready
    if _users_ready:
        return
    q = qdrant()
    try:
        q.get_collection(USERS_COL)
    except Exception:
        q.create_collection(USERS_COL, vectors_config=VectorParams(size=4, distance=Distance.DOT))
        for f in ("username", "email", "plan", "active"):
            try:
                q.create_payload_index(USERS_COL, f, PayloadSchemaType.KEYWORD)
            except Exception:
                pass
    _users_ready = True

def _ensure_logs() -> None:
    global _logs_ready
    if _logs_ready:
        return
    q = qdrant()
    try:
        q.get_collection(LOGS_COL)
    except Exception:
        q.create_collection(LOGS_COL, vectors_config=VectorParams(size=1, distance=Distance.DOT))
    _logs_ready = True

def _ensure_resets() -> None:
    global _resets_ready
    if _resets_ready:
        return
    q = qdrant()
    try:
        q.get_collection(RESETS_COL)
    except Exception:
        q.create_collection(RESETS_COL, vectors_config=VectorParams(size=4, distance=Distance.DOT))
        try:
            q.create_payload_index(RESETS_COL, "username", PayloadSchemaType.KEYWORD)
            q.create_payload_index(RESETS_COL, "token", PayloadSchemaType.KEYWORD)
        except Exception:
            pass
    _resets_ready = True

def _uid(username: str) -> str:
    return str(uuid.uuid5(uuid.NAMESPACE_DNS, f"dalilak_user_{username}"))

def db_get_user(username: str) -> Optional[dict]:
    _ensure_users()
    results, _ = qdrant().scroll(
        collection_name=USERS_COL,
        scroll_filter=Filter(must=[FieldCondition(key="username", match=MatchValue(value=username))]),
        limit=1, with_payload=True,
    )
    return results[0].payload if results else None

def db_get_user_by_email(email: str) -> Optional[dict]:
    _ensure_users()
    results, _ = qdrant().scroll(
        collection_name=USERS_COL,
        scroll_filter=Filter(must=[FieldCondition(key="email", match=MatchValue(value=email.lower()))]),
        limit=1, with_payload=True,
    )
    return results[0].payload if results else None

def db_save_user(data: dict) -> None:
    _ensure_users()
    qdrant().upsert(
        collection_name=USERS_COL,
        points=[PointStruct(id=_uid(data["username"]), vector=[0.0] * 4, payload=data)],
    )

def db_list_users() -> list[dict]:
    _ensure_users()
    results, _ = qdrant().scroll(collection_name=USERS_COL, limit=500, with_payload=True)
    return [r.payload for r in results if r.payload]

def db_save_reset(username: str, token: str, expires_at: str) -> None:
    _ensure_resets()
    rid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"reset_{username}"))
    qdrant().upsert(
        collection_name=RESETS_COL,
        points=[PointStruct(id=rid, vector=[0.0] * 4,
                            payload={"username": username, "token": token, "expires_at": expires_at, "used": False})],
    )

def db_get_reset(token: str) -> Optional[dict]:
    _ensure_resets()
    results, _ = qdrant().scroll(
        collection_name=RESETS_COL,
        scroll_filter=Filter(must=[FieldCondition(key="token", match=MatchValue(value=token))]),
        limit=1, with_payload=True,
    )
    return results[0].payload if results else None

def db_mark_reset_used(username: str) -> None:
    _ensure_resets()
    rid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"reset_{username}"))
    try:
        qdrant().set_payload(collection_name=RESETS_COL, payload={"used": True}, points=[rid])
    except Exception:
        pass

def log_query(username: str, query_type: str, elapsed_ms: int) -> None:
    try:
        _ensure_logs()
        qdrant().upsert(
            collection_name=LOGS_COL,
            points=[PointStruct(
                id=str(uuid.uuid5(uuid.NAMESPACE_DNS, str(uuid.uuid4()))),
                vector=[1.0],
                payload={
                    "username": username,
                    "type": query_type,
                    "elapsed_ms": elapsed_ms,
                    "timestamp": datetime.now(timezone.utc).isoformat(),
                },
            )],
        )
    except Exception:
        pass

# ═══════════════════════════════════════════════════════════════
#  AUTH MIDDLEWARE
# ═══════════════════════════════════════════════════════════════

_bearer = HTTPBearer(auto_error=False)

def _check_subscription(user: dict) -> None:
    """Raise 403 if trial expired and not paid."""
    plan = user.get("plan", "trial")
    if plan in ("paid", "admin", "guest"):
        return
    # trial: check expiry
    expires = user.get("trial_expires_at", "")
    if expires:
        try:
            exp_dt = datetime.fromisoformat(expires).replace(tzinfo=timezone.utc)
            if datetime.now(timezone.utc) > exp_dt:
                raise HTTPException(
                    status_code=402,
                    detail="انتهت الفترة التجريبية — يرجى الترقية إلى الاشتراك المدفوع",
                )
        except HTTPException:
            raise
        except Exception:
            pass

async def get_current_user(
    creds: Optional[HTTPAuthorizationCredentials] = Depends(_bearer),
) -> dict:
    # ── GUEST ACCESS: no token required ──────────────────────────
    if not creds:
        return {"username": "guest", "plan": "guest", "active": True, "trial_expires_at": None}
    # ── AUTHENTICATED USER ────────────────────────────────────────
    try:
        payload = decode_token(creds.credentials)
        username = payload["sub"]
        # Phase 8: check token blocklist
        jti = payload.get("jti")
        if jti and (jti in _revoked_tokens or await _jti_is_revoked_redis(jti)):
            raise HTTPException(401, detail="جلسة منتهية — سجّل الدخول مجدداً")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(401, detail="جلسة منتهية — سجّل الدخول مجدداً")
    user = db_get_user(username)
    if not user or not user.get("active", True):
        raise HTTPException(401, detail="الحساب غير موجود أو معطّل")
    _check_subscription(user)
    return user

async def get_admin_user(user: dict = Depends(get_current_user)) -> dict:
    if user.get("plan") != "admin" and user.get("role") != "admin":
        raise HTTPException(403, detail="صلاحية المشرف مطلوبة")
    return user

def verify_admin_secret(x_admin_secret: Optional[str] = Header(None)) -> None:
    if x_admin_secret != ADMIN_SECRET:
        raise HTTPException(403, detail="ADMIN_SECRET غير صحيح")

# ═══════════════════════════════════════════════════════════════
#  PYDANTIC MODELS
# ═══════════════════════════════════════════════════════════════

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    history: list[Message] = []
    domain: Optional[str] = None

class AnalyzeRequest(BaseModel):
    file_base64: str
    file_type: str
    file_name: str
    message: str = "حلل هذه الوثيقة واقترح الإجراءات المناسبة"
    history: list[Message] = []

class RegisterRequest(BaseModel):
    username: str
    email: str
    password: str
    full_name: str = ""
    phone: str = ""

class LoginRequest(BaseModel):
    username: str
    password: str

class ForgotPasswordRequest(BaseModel):
    email: str

class ResetPasswordRequest(BaseModel):
    token: str
    new_password: str

class UpdateUserRequest(BaseModel):
    plan: Optional[str] = None
    active: Optional[bool] = None
    full_name: Optional[str] = None
    phone: Optional[str] = None
    paid_until: Optional[str] = None

class CreateUserRequest(BaseModel):
    username: str
    email: str
    password: str
    full_name: str = ""
    phone: str = ""
    plan: str = "trial"

class FileExtractRequest(BaseModel):
    file_base64: str
    file_type: str
    file_name: str

# ═══════════════════════════════════════════════════════════════
#  FASTAPI APP
# ═══════════════════════════════════════════════════════════════

app = FastAPI(title="Dalilak AI API", version="4.1.0")

app.add_middleware(
    CORSMiddleware,
    # Phase 8: restrict to known origins (wildcard removed)
    allow_origins=[APP_BASE_URL, "http://localhost:3000"],
    allow_methods=["*"],
    allow_headers=["*"],
    allow_credentials=True,
)

# ═══════════════════════════════════════════════════════════════
#  RAG HELPERS
# ═══════════════════════════════════════════════════════════════

# ── Phase 9: Request-ID middleware ────────────────────────────────────────────────────────────────────────────────
# Attaches X-Request-ID to every response for client-side correlation.
# Logs method + path + status + elapsed_ms — NEVER logs query text or tokens.
@app.middleware("http")
async def _request_id_middleware(request: Request, call_next):
    rid = request.headers.get("X-Request-ID") or secrets.token_hex(8)
    _req_id_var.set(rid)
    t0 = time.time()
    response = await call_next(request)
    elapsed_ms = round((time.time() - t0) * 1000)
    response.headers["X-Request-ID"] = rid
    _log.info(
        "rid=%s method=%s path=%s status=%s elapsed_ms=%s",
        rid, request.method, request.url.path, response.status_code, elapsed_ms,
    )
    return response

# ── Phase 9: Global exception handler — no stack traces to client ─────────────────────────────────────────────────────────────────────────
@app.exception_handler(Exception)
async def _global_exception_handler(request: Request, exc: Exception):
    rid = _req_id_var.get("-")
    _log.exception("unhandled exception rid=%s: %s", rid, type(exc).__name__)
    return JSONResponse(
        status_code=500,
        content={
            "detail": "خطأ داخلي — حاول مجدداً أو تواصل مع الدعم الفني",
            "req_id": rid,
        },
    )

async def expand_query(query: str) -> str:
    """
    Lightweight query expansion: rewrite the user query to be more specific
    for semantic search. Uses gpt-4o-mini to keep latency low (~200ms).
    Only expands queries longer than 10 chars; returns original otherwise.
    """
    if len(query.strip()) <= 10:
        return query
    try:
        resp = await asyncio.wait_for(
            oai().chat.completions.create(
                model=MODEL_FAST,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "أنت محرك بحث متخصص في الإجراءات الحكومية اللبنانية. "
                            "أعد صياغة السؤال التالي في جملة بحثية واحدة أكثر دقةً وشمولاً، "
                            "مُدرِجاً المصطلحات الرسمية المحتملة. أرجع الجملة فقط."
                        ),
                    },
                    {"role": "user", "content": query},
                ],
                max_tokens=150,
                temperature=0.0,
            ),
            timeout=4.0,
        )
        expanded = resp.choices[0].message.content.strip()
        return expanded if expanded else query
    except Exception:
        return query  # Graceful degradation: use original query on any failure

async def embed(text: str) -> list:
    r = await oai().embeddings.create(
        model=EMBED_MODEL, input=[text[:MAX_CHARS]], dimensions=VECTOR_DIM,
    )
    return r.data[0].embedding

async def search_qdrant(vec: list, domain: Optional[str] = None) -> list:
    body: dict = {
        "vector": vec, "limit": MAX_CTX,
        "score_threshold": MIN_SCORE, "with_payload": True,
    }
    if domain:
        body["filter"] = {"must": [{"key": "domain", "match": {"value": domain}}]}
    async with httpx.AsyncClient(timeout=20) as client:
        r = await client.post(
            f"{qdrant_url()}/collections/{COLLECTION}/points/search",
            headers=qdrant_headers(), json=body,
        )
    items = r.json().get("result", [])
    return [
        {
            "score":    round(x.get("score", 0), 3),
            "title":    x["payload"].get("title", ""),
            "text":     x["payload"].get("text", ""),
            "domain":   x["payload"].get("domain", ""),
            "ministry": x["payload"].get("ministry", ""),
            "website":  x["payload"].get("website", ""),
            "phone":    x["payload"].get("phone", ""),
            "fees":     x["payload"].get("fees", ""),
        }
        for x in items
    ]

def _mmr_deduplicate(chunks: list, k: int = 8, diversity: float = 0.3) -> list:
    """
    Maximal Marginal Relevance: keeps top-k chunks maximising relevance while
    penalising chunks whose *title* duplicates an already-selected chunk.
    diversity ∈ [0,1] — higher = more diverse, lower = more relevant.
    """
    if len(chunks) <= k:
        return chunks
    selected: list = []
    remaining = list(chunks)
    # Always take the best-scoring chunk first
    best = max(remaining, key=lambda c: c["score"])
    selected.append(best)
    remaining.remove(best)
    while len(selected) < k and remaining:
        best_candidate, best_score = None, float("-inf")
        for cand in remaining:
            relevance = cand["score"]
            # Penalty: 1.0 if same title as any selected, else 0
            max_sim = max(
                1.0 if sel["title"] and sel["title"] == cand["title"] else 0.0
                for sel in selected
            )
            mmr = (1 - diversity) * relevance - diversity * max_sim
            if mmr > best_score:
                best_score = mmr
                best_candidate = cand
        if best_candidate:
            selected.append(best_candidate)
            remaining.remove(best_candidate)
    return selected


def context_str(chunks: list) -> str:
    if not chunks:
        return ""
    # MMR deduplication for diverse, non-redundant context
    chunks = _mmr_deduplicate(chunks, k=8)
    parts = [
        "=== المعلومات الرسمية المتاحة ===",
        "استخدم أرقام المصادر [1] [2] ... في إجابتك للإشارة إلى المصدر.",
    ]
    for i, c in enumerate(chunks, 1):
        header = f"\n[{i}] {c['title']}"
        if c["ministry"]: header += f" — {c['ministry']}"
        parts.append(header)
        parts.append(c["text"])
        extras = []
        if c.get("website"): extras.append(f"🌐 {c['website']}")
        if c.get("phone"):   extras.append(f"📞 {c['phone']}")
        if c.get("fees"):    extras.append(f"💰 الرسوم: {c['fees']}")
        if extras: parts.append("  ".join(extras))
        parts.append(f"  (درجة التطابق: {c['score']:.0%})")
        parts.append("─" * 40)
    return "\n".join(parts)

def pick_model(msg: str) -> str:
    """Use fast model only for very short greetings; always smart for real queries."""
    greetings = {"مرحبا", "اهلا", "هلا", "hi", "hello", "مساء", "صباح", "شكرا", "شكراً", "تمام", "ok", "حسنا"}
    msg_stripped = msg.strip().lower().rstrip("!.,?؟")
    if len(msg) < 30 and msg_stripped in greetings:
        return MODEL_FAST
    return MODEL_SMART

def build_messages(ctx: str, history: list, user_msg: str) -> list:
    system = SYSTEM_PROMPT + (f"\n\n{ctx}" if ctx else "")
    msgs = [{"role": "system", "content": system}]
    for m in history[-MAX_HISTORY:]:
        msgs.append({"role": m.role, "content": m.content})
    msgs.append({"role": "user", "content": user_msg})
    return msgs

# ── Evidence Gate v2 — multi-factor ──────────────────────────────────────────

class EvidenceOutcome:
    SUFFICIENT   = "SUFFICIENT"
    PARTIAL      = "PARTIAL"
    INSUFFICIENT = "INSUFFICIENT"
    CONFLICTING  = "CONFLICTING"

def _evaluate_evidence(chunks: list, query: str = "") -> tuple:
    """
    Multi-factor evidence evaluation.
    Returns (outcome, reason) where outcome is one of EvidenceOutcome constants.

    Rules:
    - INSUFFICIENT: no chunk >= SUFFICIENCY_TOP_SCORE → refuse to call LLM
    - PARTIAL: best chunk >= 0.35 but coverage < 0.45 → answer with caveats
    - SUFFICIENT: best chunk >= 0.45 and multiple chunks agree → full answer
    - CONFLICTING: top chunks have score >= 0.35 but contradicting country payloads
    """
    if not chunks:
        return EvidenceOutcome.INSUFFICIENT, "no_chunks"

    scores = [c.get("score", 0) for c in chunks]
    top_score = max(scores)

    if top_score < SUFFICIENCY_TOP_SCORE:
        return EvidenceOutcome.INSUFFICIENT, f"top_score={top_score:.3f}"

    # Check for cross-country contamination in top chunks
    top_chunks = [c for c in chunks if c.get("score", 0) >= SUFFICIENCY_TOP_SCORE]
    countries = set()
    for c in top_chunks:
        country = c.get("country", "") or c.get("domain", "")
        if country:
            countries.add(country.lower())

    # If both Lebanon and Syria appear in top results — potential conflict
    has_lb = any(c in countries for c in ["لبنان", "lebanon", "lb"])
    has_sy = any(c in countries for c in ["سوريا", "syria", "sy"])
    if has_lb and has_sy and len(top_chunks) >= 2:
        return EvidenceOutcome.CONFLICTING, "cross_country"

    # Grade coverage
    high_quality = [s for s in scores if s >= 0.45]
    if len(high_quality) >= 2:
        return EvidenceOutcome.SUFFICIENT, f"top={top_score:.3f},hq={len(high_quality)}"
    elif top_score >= 0.45:
        return EvidenceOutcome.SUFFICIENT, f"top={top_score:.3f}"
    else:
        return EvidenceOutcome.PARTIAL, f"top={top_score:.3f},partial_coverage"

# Keep backward compat — Phase 5 tests rely on this symbol being present
def _is_evidence_sufficient(chunks: list) -> bool:
    outcome, _ = _evaluate_evidence(chunks)
    return outcome != EvidenceOutcome.INSUFFICIENT

SUFFICIENCY_MSG = (
    "لم أجد في قاعدة بيانات دليلك معلومات كافية تُغطّي هذا السؤال تحديداً. "
    "أنصحك بالتواصل مع الجهة الحكومية المختصة مباشرةً للحصول على إجابة دقيقة."
)

PARTIAL_PREFIX = (
    "تنبيه: المعلومات التالية جزئية وقد لا تغطي جميع جوانب سؤالك. "
    "يُرجى التحقق من الجهة المختصة للتأكد.\n\n"
)

CONFLICT_MSG = (
    "وجدت معلومات متضاربة بين مصادر مختلفة لهذا السؤال. "
    "يُرجى تحديد البلد (لبنان أم سوريا) وإعادة السؤال للحصول على إجابة دقيقة."
)
# ──────────────────────────────────────────────────────────────────────────────

# ═══════════════════════════════════════════════════════════════
#  PUBLIC ENDPOINTS
# ═══════════════════════════════════════════════════════════════

@app.get("/")
async def root():
    return {"status": "ok", "name": "Dalilak AI", "version": "4.1.0"}

@app.get("/health")
async def health():
    report: dict = {}

    # ── Qdrant ──────────────────────────────────────────────────────────────
    qdrant_ok = False
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(f"{qdrant_url()}/collections/{COLLECTION}", headers=qdrant_headers())
        pts = r.json().get("result", {}).get("points_count", 0)
        report["qdrant"] = {"status": "ok", "collection": COLLECTION, "points": pts}
        qdrant_ok = True
    except Exception as e:
        _log.warning("health check — Qdrant failed: %s", e)
        report["qdrant"] = {"status": "unavailable"}

    # ── Redis ────────────────────────────────────────────────────────────────
    if not _REDIS_CONFIGURED_AT_STARTUP:
        report["redis"] = {
            "status": _REDIS_NOT_CONFIGURED,
            "note": "in-memory fallback active; logout revocation not durable across restarts",
        }
    else:
        report["redis"] = {"status": _redis_jti_status}

    # ── Email provider ───────────────────────────────────────────────────────
    resend_key_present = bool(os.environ.get("RESEND_API_KEY", "").strip())
    report["email"] = {
        "provider": "Resend",
        "ready": resend_key_present,
        "note": "RESEND_API_KEY not configured" if not resend_key_present else "configured",
    }

    # ── Overall status ───────────────────────────────────────────────────────
    if not qdrant_ok:
        report["status"] = "degraded"
        raise HTTPException(503, detail=report)
    report["status"] = "ok"
    return report

@app.get("/ping")
async def ping():
    return {"pong": True}

# ═══════════════════════════════════════════════════════════════
#  AUTH ENDPOINTS
# ═══════════════════════════════════════════════════════════════

@app.post("/auth/register")
async def register(req: RegisterRequest, request: Request):
    await _rate_enforce(request, "register")
    if len(req.username) < 3:
        raise HTTPException(400, detail="اسم المستخدم يجب أن يكون 3 أحرف على الأقل")
    if len(req.password) < 8:  # Phase 8: raised from 6 to 8
        raise HTTPException(400, detail="كلمة المرور يجب أن تكون 8 أحرف على الأقل")
    if "@" not in req.email:
        raise HTTPException(400, detail="البريد الإلكتروني غير صالح")
    if db_get_user(req.username.lower()):
        raise HTTPException(409, detail="اسم المستخدم محجوز — اختر اسماً آخر")
    if db_get_user_by_email(req.email.lower()):
        raise HTTPException(409, detail="البريد الإلكتروني مسجّل مسبقاً")

    now = datetime.now(timezone.utc)
    trial_expires = (now + timedelta(days=TRIAL_DAYS)).isoformat()
    user = {
        "username":         req.username.lower(),
        "email":            req.email.lower(),
        "password_hash":    hash_pw(req.password),
        "full_name":        req.full_name,
        "phone":            req.phone,
        "plan":             "trial",
        "role":             "user",
        "active":           True,
        "trial_expires_at": trial_expires,
        "paid_until":       None,
        "created_at":       now.isoformat(),
        "last_login":       None,
    }
    db_save_user(user)
    token = create_token(req.username.lower())
    return {
        "token": token,
        "user": {
            "username":         user["username"],
            "email":            user["email"],
            "full_name":        user["full_name"],
            "plan":             user["plan"],
            "trial_expires_at": trial_expires,
        },
        "message": f"مرحباً! لديك {TRIAL_DAYS} أيام تجريبية مجانية.",
    }

@app.post("/auth/login")
async def login(req: LoginRequest, request: Request):
    await _rate_enforce(request, "login")
    user = db_get_user(req.username.lower())
    if not user:
        user = db_get_user_by_email(req.username.lower())
    if not user:
        raise HTTPException(401, detail="اسم المستخدم أو كلمة المرور غير صحيحة")
    if not user.get("active", True):
        raise HTTPException(403, detail="الحساب معطّل — تواصل مع الدعم")
    if not verify_pw(req.password, user.get("password_hash", "")):
        raise HTTPException(401, detail="اسم المستخدم أو كلمة المرور غير صحيحة")

    user["last_login"] = datetime.now(timezone.utc).isoformat()
    db_save_user(user)
    token = create_token(user["username"], user.get("role", "user"))

    plan = user.get("plan", "trial")
    subscription_status = plan
    days_left = None
    if plan == "trial":
        try:
            exp = datetime.fromisoformat(user.get("trial_expires_at", "")).replace(tzinfo=timezone.utc)
            delta = (exp - datetime.now(timezone.utc)).days
            days_left = max(0, delta)
            if days_left == 0:
                subscription_status = "expired"
        except Exception:
            pass

    return {
        "token": token,
        "user": {
            "username":            user["username"],
            "email":               user["email"],
            "full_name":           user.get("full_name", ""),
            "plan":                plan,
            "role":                user.get("role", "user"),
            "trial_expires_at":    user.get("trial_expires_at"),
            "paid_until":          user.get("paid_until"),
            "subscription_status": subscription_status,
            "days_left":           days_left,
        },
    }

@app.get("/auth/me")
async def me(user: dict = Depends(get_current_user)):
    plan = user.get("plan", "trial")
    days_left = None
    if plan == "trial":
        try:
            exp = datetime.fromisoformat(user.get("trial_expires_at", "")).replace(tzinfo=timezone.utc)
            days_left = max(0, (exp - datetime.now(timezone.utc)).days)
        except Exception:
            pass
    return {
        "username":         user["username"],
        "email":            user.get("email", ""),
        "full_name":        user.get("full_name", ""),
        "plan":             plan,
        "role":             user.get("role", "user"),
        "trial_expires_at": user.get("trial_expires_at"),
        "paid_until":       user.get("paid_until"),
        "days_left":        days_left,
        "created_at":       user.get("created_at"),
    }

@app.post("/auth/forgot-password")
async def forgot_password(req: ForgotPasswordRequest, request: Request):
    await _rate_enforce(request, "forgot")
    # Anti-enumeration: always return the same message regardless of whether
    # the email exists. Never reveal whether an address is registered.
    _SAFE_RESPONSE = {"message": "إذا كان البريد مسجّلاً، ستصلك رسالة إعادة التعيين خلال دقائق."}

    user = db_get_user_by_email(req.email.lower())
    if not user:
        return _SAFE_RESPONSE

    # 32-byte URL-safe token (43 chars). Only its SHA-256 hash is stored.
    raw_token = secrets.token_urlsafe(32)
    token_hash = _hash_reset_token(raw_token)
    expires = (datetime.now(timezone.utc) + timedelta(minutes=15)).isoformat()
    db_save_reset(user["username"], token_hash, expires)

    reset_url = f"{APP_BASE_URL}/reset-password?token={raw_token}"
    email_ok = await _send_reset_email(req.email.lower(), reset_url, from_email=RESEND_FROM_EMAIL)

    if not email_ok:
        # Email delivery failed — immediately invalidate the token so it
        # cannot be used even though no email reached the user.
        # Correlation ID in log: omit token value entirely (no PII).
        corr_id = secrets.token_hex(8)
        _log.error(
            "[forgot_password][ALERT] Email delivery failed — token invalidated. "
            "corr_id=%s username_hash=%s",
            corr_id,
            _hash_reset_token(user["username"])[:8],   # partial hash, not PII
        )
        try:
            db_mark_reset_used(user["username"])
        except Exception as ex:
            _log.error("[forgot_password] Could not invalidate token: %s corr_id=%s", type(ex).__name__, corr_id)

    return _SAFE_RESPONSE

@app.post("/auth/reset-password")
async def reset_password(req: ResetPasswordRequest):
    if len(req.new_password) < 8:  # Phase 8: raised from 6 to 8
        raise HTTPException(400, detail="كلمة المرور يجب أن تكون 8 أحرف على الأقل")
    # Hash the raw token before lookup — the DB stores only hashes, never raw tokens
    token_hash = _hash_reset_token(req.token)
    reset = db_get_reset(token_hash)
    if not reset:
        raise HTTPException(400, detail="رمز الاستعادة غير صحيح")
    if reset.get("used"):
        raise HTTPException(400, detail="رمز الاستعادة مستخدم مسبقاً")
    try:
        exp = datetime.fromisoformat(reset["expires_at"]).replace(tzinfo=timezone.utc)
        if datetime.now(timezone.utc) > exp:
            raise HTTPException(400, detail="رمز الاستعادة منتهي الصلاحية")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(400, detail="رمز غير صالح")
    user = db_get_user(reset["username"])
    if not user:
        raise HTTPException(400, detail="الحساب غير موجود")
    user["password_hash"] = hash_pw(req.new_password)
    db_save_user(user)
    db_mark_reset_used(reset["username"])
    return {"message": "تم تغيير كلمة المرور بنجاح — يمكنك تسجيل الدخول الآن."}

@app.post("/auth/logout")
async def logout(
    creds: Optional[HTTPAuthorizationCredentials] = Depends(_bearer),
):
    """Phase 8: revoke the bearer token by adding its jti to the in-memory blocklist."""
    if creds:
        try:
            payload = decode_token(creds.credentials)
            jti = payload.get("jti")
            exp = payload.get("exp", time.time() + 86400)
            if jti:
                _blocklist_prune()                # evict expired entries first
                _revoked_tokens[jti] = float(exp)
                await _jti_revoke_redis(jti, float(exp))  # Phase B: persist across restarts
        except Exception:
            pass  # invalid token — treat as already logged out
    return {"message": "تم تسجيل الخروج بنجاح"}

# ═══════════════════════════════════════════════════════════════
#  ADMIN ENDPOINTS
# ═══════════════════════════════════════════════════════════════

@app.get("/admin/users")
async def admin_list_users(user: dict = Depends(get_admin_user)):
    users = db_list_users()
    now = datetime.now(timezone.utc)
    result = []
    for u in users:
        plan = u.get("plan", "trial")
        status = plan
        days_left = None
        if plan == "trial":
            try:
                exp = datetime.fromisoformat(u.get("trial_expires_at", "")).replace(tzinfo=timezone.utc)
                dl = (exp - now).days
                days_left = max(0, dl)
                if dl < 0:
                    status = "expired"
            except Exception:
                pass
        result.append({
            "username":   u.get("username"),
            "email":      u.get("email"),
            "full_name":  u.get("full_name", ""),
            "phone":      u.get("phone", ""),
            "plan":       plan,
            "status":     status,
            "days_left":  days_left,
            "active":     u.get("active", True),
            "created_at": u.get("created_at"),
            "last_login": u.get("last_login"),
            "paid_until": u.get("paid_until"),
        })
    result.sort(key=lambda x: x.get("created_at") or "", reverse=True)
    return {"users": result, "total": len(result)}

@app.post("/admin/users")
async def admin_create_user(req: CreateUserRequest, user: dict = Depends(get_admin_user)):
    if db_get_user(req.username.lower()):
        raise HTTPException(409, detail="اسم المستخدم محجوز")
    if db_get_user_by_email(req.email.lower()):
        raise HTTPException(409, detail="البريد مسجّل مسبقاً")
    now = datetime.now(timezone.utc)
    trial_expires = (now + timedelta(days=TRIAL_DAYS)).isoformat()
    new_user = {
        "username":         req.username.lower(),
        "email":            req.email.lower(),
        "password_hash":    hash_pw(req.password),
        "full_name":        req.full_name,
        "phone":            req.phone,
        "plan":             req.plan,
        "role":             "admin" if req.plan == "admin" else "user",
        "active":           True,
        "trial_expires_at": trial_expires,
        "paid_until":       None,
        "created_at":       now.isoformat(),
        "last_login":       None,
    }
    db_save_user(new_user)
    return {"message": f"تم إنشاء المستخدم {req.username}", "user": new_user}

@app.put("/admin/users/{username}")
async def admin_update_user(
    username: str,
    req: UpdateUserRequest,
    admin: dict = Depends(get_admin_user),
):
    user = db_get_user(username.lower())
    if not user:
        raise HTTPException(404, detail="المستخدم غير موجود")
    if req.plan is not None:
        user["plan"] = req.plan
        if req.plan == "admin":
            user["role"] = "admin"
    if req.active is not None:
        user["active"] = req.active
    if req.full_name is not None:
        user["full_name"] = req.full_name
    if req.phone is not None:
        user["phone"] = req.phone
    if req.paid_until is not None:
        user["paid_until"] = req.paid_until
        user["plan"] = "paid"
    db_save_user(user)
    return {"message": "تم التحديث بنجاح", "user": user}

@app.delete("/admin/users/{username}")
async def admin_deactivate_user(username: str, admin: dict = Depends(get_admin_user)):
    user = db_get_user(username.lower())
    if not user:
        raise HTTPException(404, detail="المستخدم غير موجود")
    user["active"] = False
    db_save_user(user)
    return {"message": f"تم تعطيل حساب {username}"}

@app.get("/admin/stats")
async def admin_stats(admin: dict = Depends(get_admin_user)):
    users = db_list_users()
    now = datetime.now(timezone.utc)
    total = len(users)
    paid = sum(1 for u in users if u.get("plan") == "paid")
    trial_active = trial_expired = suspended = 0
    for u in users:
        plan = u.get("plan", "trial")
        if not u.get("active", True):
            suspended += 1
        elif plan == "trial":
            try:
                exp = datetime.fromisoformat(u.get("trial_expires_at", "")).replace(tzinfo=timezone.utc)
                if datetime.now(timezone.utc) > exp:
                    trial_expired += 1
                else:
                    trial_active += 1
            except Exception:
                trial_active += 1
    return {
        "total": total, "paid": paid,
        "trial_active": trial_active, "trial_expired": trial_expired,
        "suspended": suspended,
        "conversion_rate": f"{round(paid / total * 100, 1)}%" if total else "0%",
    }

@app.get("/admin/resets")
async def admin_list_resets(admin: dict = Depends(get_admin_user)):
    _ensure_resets()
    results, _ = qdrant().scroll(collection_name=RESETS_COL, limit=100, with_payload=True)
    codes = [r.payload for r in results if r.payload and not r.payload.get("used")]
    now = datetime.now(timezone.utc)
    active = []
    for c in codes:
        try:
            exp = datetime.fromisoformat(c["expires_at"]).replace(tzinfo=timezone.utc)
            if now < exp:
                active.append({
                    "username":     c["username"],
                    "expires_at":   c["expires_at"],
                    # token hash omitted from response — admin does not need it
                    # (Phase B: removing hash exposure from admin endpoint)
                })
        except Exception:
            pass
    return {"reset_codes": active}

# ── Extend trial ─────────────────────────────────────────────
class ExtendTrialRequest(BaseModel):
    days: int = 7

@app.post("/admin/users/{username}/extend-trial")
async def admin_extend_trial(username: str, req: ExtendTrialRequest, admin: dict = Depends(get_admin_user)):
    user = db_get_user(username.lower())
    if not user:
        raise HTTPException(404, detail="المستخدم غير موجود")
    now = datetime.now(timezone.utc)
    # Extend from now or from current expiry, whichever is later
    current_exp_str = user.get("trial_expires_at") or now.isoformat()
    try:
        current_exp = datetime.fromisoformat(current_exp_str).replace(tzinfo=timezone.utc)
        base = max(now, current_exp)
    except Exception:
        base = now
    new_exp = (base + timedelta(days=req.days)).isoformat()
    user["trial_expires_at"] = new_exp
    if user.get("plan") not in ("paid", "admin"):
        user["plan"] = "trial"
    user["active"] = True
    db_save_user(user)
    return {"message": f"تم تمديد المهلة {req.days} أيام", "trial_expires_at": new_exp}

# ── Per-user logs ────────────────────────────────────────────
@app.get("/admin/users/{username}/logs")
async def admin_user_logs(username: str, admin: dict = Depends(get_admin_user)):
    _ensure_logs()
    results, _ = qdrant().scroll(
        collection_name=LOGS_COL,
        scroll_filter=Filter(must=[FieldCondition(key="username", match=MatchValue(value=username.lower()))]),
        limit=200, with_payload=True,
    )
    logs = sorted(
        [r.payload for r in results if r.payload],
        key=lambda x: x.get("timestamp", ""), reverse=True,
    )
    total = len(logs)
    chat_count = sum(1 for l in logs if "chat" in l.get("type", ""))
    analyze_count = sum(1 for l in logs if "analyze" in l.get("type", ""))
    avg_ms = int(sum(l.get("elapsed_ms", 0) for l in logs) / total) if total else 0
    return {
        "username": username,
        "total_queries": total,
        "chat_count": chat_count,
        "analyze_count": analyze_count,
        "avg_response_ms": avg_ms,
        "logs": logs[:50],  # last 50
    }

# ── All logs (general report) ────────────────────────────────
@app.get("/admin/logs")
async def admin_all_logs(admin: dict = Depends(get_admin_user)):
    _ensure_logs()
    results, _ = qdrant().scroll(collection_name=LOGS_COL, limit=1000, with_payload=True)
    logs = [r.payload for r in results if r.payload]
    total = len(logs)
    chat_count = sum(1 for l in logs if "chat" in l.get("type", ""))
    analyze_count = sum(1 for l in logs if "analyze" in l.get("type", ""))
    avg_ms = int(sum(l.get("elapsed_ms", 0) for l in logs) / total) if total else 0
    # Daily activity (last 14 days)
    now = datetime.now(timezone.utc)
    daily: dict = {}
    for l in logs:
        try:
            d = datetime.fromisoformat(l.get("timestamp", "")).replace(tzinfo=timezone.utc)
            if (now - d).days <= 13:
                key = d.strftime("%Y-%m-%d")
                daily[key] = daily.get(key, 0) + 1
        except Exception:
            pass
    # Top users
    user_counts: dict = {}
    for l in logs:
        u = l.get("username", "guest")
        user_counts[u] = user_counts.get(u, 0) + 1
    top_users = sorted(user_counts.items(), key=lambda x: x[1], reverse=True)[:10]
    return {
        "total_queries": total,
        "chat_count": chat_count,
        "analyze_count": analyze_count,
        "avg_response_ms": avg_ms,
        "daily_activity": daily,
        "top_users": [{"username": u, "count": c} for u, c in top_users],
    }

# ── Knowledge base management ────────────────────────────────
class KnowledgeAddRequest(BaseModel):
    title: str
    text: str
    domain: str = ""
    ministry: str = ""
    website: str = ""
    phone: str = ""
    fees: str = ""
    source: str = "admin"

@app.post("/admin/knowledge/add")
async def admin_knowledge_add(req: KnowledgeAddRequest, admin: dict = Depends(get_admin_user)):
    if not req.title.strip() or not req.text.strip():
        raise HTTPException(400, detail="العنوان والنص مطلوبان")
    full_text = f"{req.title}\n{req.text}"
    vec = await embed(full_text)
    point_id = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"admin_{req.title}_{req.text[:50]}"))
    payload = {
        "title":    req.title.strip(),
        "text":     req.text.strip(),
        "domain":   req.domain.strip(),
        "ministry": req.ministry.strip(),
        "website":  req.website.strip(),
        "phone":    req.phone.strip(),
        "fees":     req.fees.strip(),
        "source":   req.source,
        "added_by": admin.get("username"),
        "added_at": datetime.now(timezone.utc).isoformat(),
    }
    qdrant().upsert(
        collection_name=COLLECTION,
        points=[PointStruct(id=point_id, vector=vec, payload=payload)],
    )
    return {"message": "تمت إضافة المعلومة للقاعدة بنجاح", "id": point_id}

@app.get("/admin/knowledge/search")
async def admin_knowledge_search(q: str, admin: dict = Depends(get_admin_user)):
    if not q.strip():
        raise HTTPException(400, detail="أدخل نص البحث")
    vec = await embed(q)
    chunks = await search_qdrant(vec)
    return {"query": q, "results": chunks, "count": len(chunks)}

@app.get("/admin/knowledge/count")
async def admin_knowledge_count(admin: dict = Depends(get_admin_user)):
    async with httpx.AsyncClient(timeout=10) as client:
        r = await client.get(f"{qdrant_url()}/collections/{COLLECTION}", headers=qdrant_headers())
    info = r.json().get("result", {})
    return {
        "collection": COLLECTION,
        "points_count": info.get("points_count", 0),
        "segments_count": info.get("segments_count", 0),
    }

@app.post("/admin/knowledge/extract")
async def admin_knowledge_extract(req: FileExtractRequest, admin: dict = Depends(get_admin_user)):
    """Extract structured knowledge from uploaded file (PDF, Word, Excel, image, text)."""
    fname = req.file_name.lower()
    ftype = req.file_type

    # ── Extract raw text ────────────────────────────────────────
    raw_text = ""
    try:
        if ftype == "application/pdf" or fname.endswith(".pdf"):
            raw_text = extract_text_from_pdf(req.file_base64)

        elif "word" in ftype or fname.endswith((".docx", ".doc")):
            raw_text = extract_text_from_docx(req.file_base64)

        elif fname.endswith((".xlsx", ".xls", ".csv")):
            try:
                if fname.endswith(".csv"):
                    import csv, io as _io
                    decoded = base64.b64decode(req.file_base64).decode("utf-8", errors="replace")
                    reader = csv.reader(_io.StringIO(decoded))
                    rows = [" | ".join(r) for r in reader]
                    raw_text = "\n".join(rows[:200])
                else:
                    # openpyxl for xlsx
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(base64.b64decode(req.file_base64)), data_only=True)
                    parts = []
                    for ws in wb.worksheets[:3]:
                        parts.append(f"[ورقة: {ws.title}]")
                        for row in ws.iter_rows(max_row=100, values_only=True):
                            cells = [str(c) for c in row if c is not None]
                            if cells:
                                parts.append(" | ".join(cells))
                    raw_text = "\n".join(parts)
            except Exception as ex:
                raw_text = f"[تعذّر استخراج الجدول: {ex}]"

        elif ftype.startswith("image/"):
            # Use GPT-4o vision to extract text from image
            resp = await oai().chat.completions.create(
                model=MODEL_SMART,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "استخرج كل النصوص والمعلومات من هذه الصورة بدقة."},
                        {"type": "image_url", "image_url": {"url": f"data:{ftype};base64,{req.file_base64}", "detail": "high"}},
                    ],
                }],
                max_tokens=2000,
            )
            raw_text = resp.choices[0].message.content

        elif ftype.startswith("text/") or fname.endswith(".txt"):
            raw_text = base64.b64decode(req.file_base64).decode("utf-8", errors="replace")[:15000]

        else:
            raw_text = f"[نوع الملف {ftype} غير مدعوم مباشرةً — يُرجى نسخ المحتوى يدوياً]"
    except Exception as e:
        raw_text = f"[خطأ في الاستخراج: {e}]"

    if not raw_text.strip():
        return {"error": "تعذّر استخراج نص من الملف", "raw_text": ""}

    # ── Ask GPT to structure it for KB ──────────────────────────
    structure_prompt = f"""أنت مساعد لإدارة قاعدة معرفة حكومية لبنانية.
النص المستخرج من الملف "{req.file_name}":

{raw_text[:8000]}

استخرج منه المعلومات وأعدها بتنسيق JSON صارم كالتالي (لا تضف أي نص خارج JSON):
{{
  "entries": [
    {{
      "title": "عنوان واضح ومختصر للموضوع",
      "text": "النص الكامل والمفيد للمستخدم (خطوات، شروط، وثائق مطلوبة، مواعيد، رسوم، معلومات مفيدة)",
      "domain": "القطاع (مثال: سفر، تعليم، عقارات، شركات، صحة، مركبات، قانون)",
      "ministry": "الجهة المختصة أو الوزارة",
      "website": "الموقع الإلكتروني إن وجد",
      "phone": "رقم الهاتف إن وجد",
      "fees": "الرسوم المطلوبة إن وجدت"
    }}
  ],
  "summary": "ملخص قصير لمحتوى الملف"
}}

إذا كان الملف يحتوي على عدة مواضيع مختلفة، اجعل entries متعددة (حد أقصى 5). إذا كان موضوعاً واحداً، اجعل entries واحداً فقط."""

    try:
        resp = await oai().chat.completions.create(
            model=MODEL_SMART,
            messages=[{"role": "user", "content": structure_prompt}],
            max_tokens=3000,
            temperature=0.1,
            response_format={"type": "json_object"},
        )
        import json as _json
        result = _json.loads(resp.choices[0].message.content)
        return {
            "file_name": req.file_name,
            "raw_text_preview": raw_text[:500],
            "entries": result.get("entries", []),
            "summary": result.get("summary", ""),
        }
    except Exception as e:
        # Fallback: return raw text as single entry
        return {
            "file_name": req.file_name,
            "raw_text_preview": raw_text[:500],
            "entries": [{"title": req.file_name, "text": raw_text[:3000], "domain": "", "ministry": "", "website": "", "phone": "", "fees": ""}],
            "summary": "",
        }

# ═══════════════════════════════════════════════════════════════
#  CHAT ENDPOINTS
# ═══════════════════════════════════════════════════════════════

@app.post("/chat")
async def chat(req: ChatRequest, request: Request, user: dict = Depends(get_current_user)):
    await _rate_enforce(request, "chat", user_id=user["username"])
    await _check_quota(user["username"], user.get("plan", "trial"))   # Phase 10
    # Phase 12: reject oversized messages before any AI work
    if len(req.message) > MAX_MESSAGE_LEN:
        raise HTTPException(400, detail="الرسالة طويلة جداً — الحد الأقصى 4000 حرف")
    ck = _ck(req.message, req.domain)
    cached = _cget(ck)
    if cached:
        return cached

    t0 = time.time()
    search_query = await expand_query(req.message)
    vec    = await embed(search_query)
    chunks = await search_qdrant(vec, req.domain)

    # ── Phase 5 evidence gate (multi-factor v2) ──────────────────────────────
    outcome, gate_reason = _evaluate_evidence(chunks, req.message)
    _log.info("[chat] evidence_gate outcome=%s reason=%s query_expanded=%s", outcome, gate_reason, search_query != req.message)
    if not _is_evidence_sufficient(chunks):
        return {"answer": SUFFICIENCY_MSG, "model": "gate", "chunks_used": 0,
                "elapsed_s": round(time.time() - t0, 2), "sources": [], "gate": "INSUFFICIENT"}
    if outcome == EvidenceOutcome.CONFLICTING:
        return {"answer": CONFLICT_MSG, "model": "gate", "chunks_used": 0,
                "elapsed_s": round(time.time() - t0, 2), "sources": [], "gate": "CONFLICTING"}
    # ──────────────────────────────────────────────────────────────────────────

    ctx    = context_str(chunks)
    user_msg = req.message
    if outcome == EvidenceOutcome.PARTIAL:
        user_msg = f"[تنبيه: المعلومات جزئية، أجب فقط بما تجده في المصادر] {user_msg}"
    model  = pick_model(req.message)
    msgs   = build_messages(ctx, req.history, user_msg)

    resp = await oai().chat.completions.create(
        model=model, messages=msgs, max_tokens=MAX_TOKENS, temperature=0.3,
    )
    elapsed_ms = int((time.time() - t0) * 1000)
    result = {
        "answer":      resp.choices[0].message.content,
        "model":       model,
        "chunks_used": len(chunks),
        "elapsed_s":   round(elapsed_ms / 1000, 2),
        "sources": [{"title": c["title"], "ministry": c["ministry"], "score": c["score"]} for c in chunks[:5]],
    }
    _cset(ck, result)
    log_query(user["username"], "chat", elapsed_ms)
    return result

@app.post("/chat/stream")
async def chat_stream(req: ChatRequest, request: Request, user: dict = Depends(get_current_user)):
    await _rate_enforce(request, "chat", user_id=user["username"])
    await _check_quota(user["username"], user.get("plan", "trial"))   # Phase 10
    # Phase 12: reject oversized messages before any AI work
    if len(req.message) > MAX_MESSAGE_LEN:
        raise HTTPException(400, detail="الرسالة طويلة جداً — الحد الأقصى 4000 حرف")
    async def generate() -> AsyncGenerator[str, None]:
        try:
            t0 = time.time()
            search_query = await expand_query(req.message)
            vec    = await embed(search_query)
            chunks = await search_qdrant(vec, req.domain)

            # ── Phase 5 evidence gate (multi-factor v2) ───────────────────────
            outcome, gate_reason = _evaluate_evidence(chunks, req.message)
            _log.info("[chat] evidence_gate outcome=%s reason=%s query_expanded=%s", outcome, gate_reason, search_query != req.message)
            if not _is_evidence_sufficient(chunks):
                gate_ev = {"type": "gate", "answer": SUFFICIENCY_MSG, "sources": [], "gate": "INSUFFICIENT"}
                yield f"data: {json.dumps(gate_ev, ensure_ascii=False)}\n\n"
                yield "data: [DONE]\n\n"
                return
            if outcome == EvidenceOutcome.CONFLICTING:
                gate_ev = {"type": "gate", "answer": CONFLICT_MSG, "sources": [], "gate": "CONFLICTING"}
                yield f"data: {json.dumps(gate_ev, ensure_ascii=False)}\n\n"
                yield "data: [DONE]\n\n"
                return
            # ──────────────────────────────────────────────────────────────────

            ctx    = context_str(chunks)
            user_msg = req.message
            if outcome == EvidenceOutcome.PARTIAL:
                user_msg = f"[تنبيه: المعلومات جزئية، أجب فقط بما تجده في المصادر] {user_msg}"
            model  = pick_model(req.message)
            msgs   = build_messages(ctx, req.history, user_msg)

            meta = {
                "type": "meta", "model": model, "chunks": len(chunks),
                "sources": [{"title": c["title"], "ministry": c["ministry"], "score": c["score"]} for c in chunks[:5]],
            }
            yield f"data: {json.dumps(meta, ensure_ascii=False)}\n\n"

            stream = await oai().chat.completions.create(
                model=model, messages=msgs, max_tokens=MAX_TOKENS, temperature=0.3, stream=True,
            )
            async for chunk in stream:
                delta = chunk.choices[0].delta.content
                if delta:
                    yield f"data: {json.dumps({'type':'token','text':delta,'choices':[{'delta':{'content':delta}}]}, ensure_ascii=False)}\n\n"

            yield "data: [DONE]\n\n"
            log_query(user["username"], "chat_stream", int((time.time() - t0) * 1000))
        except Exception as e:
            yield f"data: {json.dumps({'type':'error','detail':str(e)})}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

@app.post("/analyze/stream")
async def analyze_stream(req: AnalyzeRequest, user: dict = Depends(get_current_user)):
    async def generate() -> AsyncGenerator[str, None]:
        try:
            fname_lower = req.file_name.lower()
            ftype       = req.file_type or ""

            is_image = ftype.startswith("image/")
            is_pdf   = ftype == "application/pdf" or fname_lower.endswith(".pdf")
            is_word  = "word" in ftype or fname_lower.endswith((".docx", ".doc"))
            is_excel = "excel" in ftype or "spreadsheet" in ftype or fname_lower.endswith((".xlsx", ".xls"))
            is_pptx  = "presentation" in ftype or fname_lower.endswith((".pptx", ".ppt"))
            is_csv   = ftype == "text/csv" or fname_lower.endswith(".csv")
            is_audio = ftype.startswith("audio/") or fname_lower.endswith((".mp3", ".wav", ".m4a", ".ogg", ".flac", ".webm", ".aac"))
            is_zip   = "zip" in ftype or fname_lower.endswith((".zip", ".rar", ".7z"))
            is_text  = ftype.startswith("text/") or fname_lower.endswith((".txt", ".md", ".json", ".xml", ".html", ".htm", ".css", ".js", ".ts", ".py", ".java", ".cpp", ".c", ".sh"))

            extracted_text = ""
            if is_pdf:
                extracted_text = extract_text_from_pdf(req.file_base64)
            elif is_word:
                extracted_text = extract_text_from_docx(req.file_base64)
            elif is_excel:
                extracted_text = extract_text_from_excel(req.file_base64)
            elif is_pptx:
                extracted_text = extract_text_from_pptx(req.file_base64)
            elif is_csv:
                extracted_text = extract_text_from_csv(req.file_base64)
            elif is_zip:
                extracted_text = extract_text_from_zip(req.file_base64)
            elif is_audio:
                extracted_text = await transcribe_audio(req.file_base64, ftype, req.file_name)
            elif is_text:
                try:
                    extracted_text = base64.b64decode(req.file_base64).decode("utf-8", errors="replace")[:15000]
                except Exception:
                    extracted_text = ""

            search_query = f"{req.file_name} {req.message} {extracted_text[:300]}"
            try:
                vec    = await embed(search_query)
                chunks = await search_qdrant(vec)
                ctx    = context_str(chunks)
            except Exception:
                ctx = ""

            ANALYSIS_PROMPT = SYSTEM_PROMPT + """

---

## قواعد تحليل الوثائق

أنت خبير متخصص في تحليل الوثائق الرسمية والقانونية اللبنانية.
عند تحليل أي وثيقة، اتبع هذا الهيكل الإلزامي بالترتيب:

### 1. 📋 تشخيص الوثيقة
- نوعها الدقيق (عقد / قرار / طلب / فاتورة / قيد / وكالة / حكم / إلخ)
- الجهة المُصدِرة والجهة المُستلِمة
- التاريخ ورقم المرجع إن وجد

### 2. 📌 استخراج البيانات الجوهرية
استخرج كل المعلومات المهمة: أسماء، أرقام، مبالغ، مواعيد، شروط، التزامات.

### 3. ⚠️ التنبيهات والمخاطر
هل هناك مواعيد نهائية قريبة؟ بنود مُلزِمة؟ إجراءات واجبة لم تُنفَّذ؟ تناقضات؟

### 4. ✅ الإجراءات العملية المطلوبة (بالترتيب)
خطوات واضحة ومرقّمة يجب على المواطن اتخاذها.

### 5. 📁 المستندات والمتطلبات
ما يجب تحضيره: وثائق، صور، طوابع، رسوم.

### 6. 🏛️ الجهة المختصة والتواصل
الوزارة أو الدائرة المختصة، رقم الهاتف، ساعات العمل.

### 7. 📝 النموذج أو المسودة الجاهزة
**إلزامي:** إذا استوجبت الوثيقة طلباً أو إفادةً: اكتب مسودة جاهزة بصيغة رسمية.

---
""" + (f"\n\n{ctx}" if ctx else "")

            file_label = (
                "صورة" if is_image else
                "ملف PDF" if is_pdf else
                "مستند Word" if is_word else
                "ملف Excel" if is_excel else
                "عرض تقديمي PowerPoint" if is_pptx else
                "ملف CSV" if is_csv else
                "ملف صوتي" if is_audio else
                "أرشيف مضغوط" if is_zip else
                "ملف نصي/كود" if is_text else
                "ملف"
            )
            user_text = f"سؤال/طلب المستخدم: {req.message}\n\nاسم الملف: {req.file_name} ({file_label})"
            if extracted_text and not extracted_text.startswith("[تعذّر"):
                label = "النص المُستخرج" if not is_audio else "النص المُحوَّل من الصوت"
                user_text += f"\n\n--- {label} ---\n{extracted_text}\n--- نهاية ---"
            elif extracted_text.startswith("[تعذّر"):
                user_text += f"\n\nملاحظة: {extracted_text}"

            if is_image:
                user_content: list = [
                    {"type": "image_url", "image_url": {"url": f"data:{ftype};base64,{req.file_base64}", "detail": "high"}},
                    {"type": "text", "text": user_text},
                ]
            else:
                user_content = [{"type": "text", "text": user_text}]

            msgs: list = [{"role": "system", "content": ANALYSIS_PROMPT}]
            for m in req.history[-MAX_HISTORY:]:
                msgs.append({"role": m.role, "content": m.content})
            msgs.append({"role": "user", "content": user_content})

            stream = await oai().chat.completions.create(
                model=MODEL_SMART, messages=msgs, max_tokens=MAX_DOC_TOKENS,
                temperature=0.2, stream=True,
            )
            async for chunk in stream:
                delta = chunk.choices[0].delta.content
                if delta:
                    yield f"data: {json.dumps({'type':'token','text':delta,'choices':[{'delta':{'content':delta}}]}, ensure_ascii=False)}\n\n"

            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type':'error','detail':str(e)})}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

# ══════════════════════════════════════════════════════════════
#  PROCEDURES API — Lebanon government service catalog
# ══════════════════════════════════════════════════════════════

# Built-in procedure catalog (seed data - expandable via admin)
_PROCEDURE_CATALOG: list[dict] = [
    {
        "id": "passport-renewal-lb",
        "slug": "passport-renewal-lb",
        "title_ar": "تجديد جواز السفر اللبناني",
        "title_en": "Lebanese Passport Renewal",
        "category": "civil_status",
        "country": "lebanon",
        "authority": "المديرية العامة للأمن العام",
        "authority_en": "General Directorate of General Security",
        "authority_url": "https://www.general-security.gov.lb",
        "summary_ar": "تجديد جواز السفر اللبناني للمواطنين في لبنان وخارجه",
        "duration_ar": "5-10 أيام عمل",
        "fee_ar": "120,000 ل.ل. للبالغين (3 سنوات)، 180,000 ل.ل. (5 سنوات)",
        "documents": [
            {"name_ar": "جواز السفر القديم", "required": True},
            {"name_ar": "هوية شخصية سارية", "required": True},
            {"name_ar": "صورة شخصية حديثة 4x4", "required": True},
            {"name_ar": "وثيقة عائلية (للمتزوجين)", "required": False},
        ],
        "steps": [
            {"order": 1, "title_ar": "تقديم الطلب", "desc_ar": "تقديم الطلب عبر الموقع الإلكتروني أو مباشرةً في مراكز الأمن العام"},
            {"order": 2, "title_ar": "دفع الرسوم", "desc_ar": "دفع الرسوم عبر OMT أو في المركز"},
            {"order": 3, "title_ar": "البصمة والصورة", "desc_ar": "أخذ البصمة والصورة في المركز"},
            {"order": 4, "title_ar": "الاستلام", "desc_ar": "استلام جواز السفر المجدَّد"},
        ],
        "source_tier": 2,
        "last_verified": "2024-01",
        "review_expiry": "2025-01",
        "status": "verified",
        "tags": ["جواز", "سفر", "أمن عام", "هوية"],
    },
    {
        "id": "birth-registration-lb",
        "slug": "birth-registration-lb",
        "title_ar": "تسجيل عقد الولادة",
        "title_en": "Birth Registration",
        "category": "civil_status",
        "country": "lebanon",
        "authority": "وزارة الداخلية والبلديات — دوائر النفوس",
        "authority_en": "Ministry of Interior — Civil Registry",
        "summary_ar": "تسجيل المواليد الجدد في السجلات المدنية اللبنانية",
        "duration_ar": "يوم واحد (إذا كانت الوثائق مكتملة)",
        "fee_ar": "مجاني (خلال 30 يوماً من الولادة)",
        "documents": [
            {"name_ar": "شهادة الولادة من المستشفى", "required": True},
            {"name_ar": "هويات الوالدين", "required": True},
            {"name_ar": "عقد زواج الوالدين", "required": True},
        ],
        "steps": [
            {"order": 1, "title_ar": "الحصول على شهادة الولادة", "desc_ar": "من المستشفى خلال 24 ساعة"},
            {"order": 2, "title_ar": "التوجه لدائرة النفوس", "desc_ar": "في قضاء إقامة الأب"},
            {"order": 3, "title_ar": "تقديم الوثائق", "desc_ar": "تقديم كافة الوثائق المطلوبة"},
            {"order": 4, "title_ar": "استلام عقد الولادة", "desc_ar": "في نفس اليوم عادةً"},
        ],
        "source_tier": 2,
        "last_verified": "2024-01",
        "review_expiry": "2025-01",
        "status": "verified",
        "tags": ["ولادة", "تسجيل", "نفوس", "مواليد"],
    },
    {
        "id": "vehicle-registration-lb",
        "slug": "vehicle-registration-lb",
        "title_ar": "تسجيل سيارة جديدة",
        "title_en": "New Vehicle Registration",
        "category": "vehicles",
        "country": "lebanon",
        "authority": "مديرية السير العام",
        "authority_en": "General Traffic Directorate",
        "summary_ar": "تسجيل مركبة جديدة أو نقل ملكية مركبة في لبنان",
        "duration_ar": "1-3 أيام عمل",
        "fee_ar": "تعتمد على نوع المركبة والمحرك",
        "documents": [
            {"name_ar": "فاتورة الشراء", "required": True},
            {"name_ar": "وثيقة الاستيراد (للسيارات المستوردة)", "required": True},
            {"name_ar": "شهادة الفحص التقني", "required": True},
            {"name_ar": "وثيقة التأمين", "required": True},
            {"name_ar": "هوية المالك", "required": True},
        ],
        "steps": [
            {"order": 1, "title_ar": "الفحص التقني", "desc_ar": "إجراء الفحص التقني للمركبة"},
            {"order": 2, "title_ar": "التأمين", "desc_ar": "الحصول على وثيقة تأمين سارية"},
            {"order": 3, "title_ar": "تقديم الأوراق", "desc_ar": "تقديم جميع الوثائق لمديرية السير"},
            {"order": 4, "title_ar": "دفع الرسوم", "desc_ar": "دفع رسوم التسجيل"},
            {"order": 5, "title_ar": "استلام اللوحة والدفتر", "desc_ar": "استلام لوحة السير ودفتر السيارة"},
        ],
        "source_tier": 2,
        "last_verified": "2024-01",
        "review_expiry": "2025-01",
        "status": "verified",
        "tags": ["سيارة", "مركبة", "تسجيل", "سير", "لوحة"],
    },
    {
        "id": "company-registration-lb",
        "slug": "company-registration-lb",
        "title_ar": "تأسيس شركة في لبنان",
        "title_en": "Company Registration in Lebanon",
        "category": "business",
        "country": "lebanon",
        "authority": "وزارة الاقتصاد والتجارة — الإدارة المركزية للإحصاء",
        "authority_en": "Ministry of Economy and Trade",
        "summary_ar": "تسجيل وتأسيس شركة تجارية أو شركة ش.م.م. في لبنان",
        "duration_ar": "15-30 يوم عمل",
        "fee_ar": "تعتمد على نوع الشركة ورأس المال",
        "documents": [
            {"name_ar": "عقد التأسيس موثق عند كاتب عدل", "required": True},
            {"name_ar": "هويات المؤسسين", "required": True},
            {"name_ar": "إيصال دفع رأس المال", "required": True},
            {"name_ar": "إيجار المقر أو وثيقة الملكية", "required": True},
        ],
        "steps": [
            {"order": 1, "title_ar": "اختيار اسم الشركة", "desc_ar": "التحقق من توفر الاسم في السجل التجاري"},
            {"order": 2, "title_ar": "توثيق عقد التأسيس", "desc_ar": "لدى كاتب عدل معتمد"},
            {"order": 3, "title_ar": "فتح حساب بنكي", "desc_ar": "إيداع رأس المال في البنك"},
            {"order": 4, "title_ar": "التسجيل في السجل التجاري", "desc_ar": "وزارة العدل"},
            {"order": 5, "title_ar": "الحصول على السجل الضريبي", "desc_ar": "وزارة المالية"},
            {"order": 6, "title_ar": "التسجيل في الضمان الاجتماعي", "desc_ar": "إذا كان هناك موظفون"},
        ],
        "source_tier": 2,
        "last_verified": "2024-01",
        "review_expiry": "2025-01",
        "status": "verified",
        "tags": ["شركة", "تأسيس", "تجارة", "سجل تجاري", "ش.م.م."],
    },
    {
        "id": "driving-license-lb",
        "slug": "driving-license-lb",
        "title_ar": "استخراج رخصة قيادة جديدة",
        "title_en": "New Driving License",
        "category": "vehicles",
        "country": "lebanon",
        "authority": "مديرية السير والمركبات الآلية",
        "authority_en": "Traffic and Motor Vehicles Directorate",
        "authority_url": "https://www.imta.gov.lb",
        "summary_ar": "استخراج رخصة قيادة لأول مرة للمقيمين في لبنان بعد اجتياز الاختبارات النظرية والعملية.",
        "duration_ar": "4-8 أسابيع",
        "fee_ar": "تقريباً 150,000 ل.ل.",
        "tags": ["سيارة", "رخصة", "قيادة", "سير"],
        "source_tier": 2,
        "last_verified": "2024-01-01",
        "review_expiry": "2025-01-01",
        "status": "verified",
        "documents": [
            {"name_ar": "بطاقة الهوية أو جواز السفر", "required": True},
            {"name_ar": "شهادة طبية من طبيب مرخص", "required": True},
            {"name_ar": "صورتان شمسيتان", "required": True},
            {"name_ar": "إيصال دفع الرسوم", "required": True},
        ],
        "steps": [
            {"order": 1, "title_ar": "التسجيل في مدرسة تعليم القيادة", "desc_ar": "التسجيل في إحدى مدارس القيادة المعتمدة وأخذ 20 حصة نظرية و10 عملية على الأقل."},
            {"order": 2, "title_ar": "اجتياز اختبار الكود", "desc_ar": "التقدم لاختبار الكود النظري (30 سؤالاً) لدى مديرية السير."},
            {"order": 3, "title_ar": "اجتياز الاختبار العملي", "desc_ar": "التقدم للاختبار العملي على الطريق مع محقق السير."},
            {"order": 4, "title_ar": "تقديم الملف وسداد الرسوم", "desc_ar": "تقديم كامل الوثائق ودفع رسوم استخراج الرخصة."},
            {"order": 5, "title_ar": "استلام الرخصة", "desc_ar": "استلام رخصة القيادة خلال أسبوع إلى أسبوعين."},
        ],
    },
    {
        "id": "marriage-certificate-lb",
        "slug": "marriage-certificate-lb",
        "title_ar": "استخراج عقد الزواج المدني",
        "title_en": "Marriage Certificate",
        "category": "civil_status",
        "country": "lebanon",
        "authority": "وزارة الداخلية والبلديات",
        "authority_en": "Ministry of Interior and Municipalities",
        "authority_url": "https://www.interior.gov.lb",
        "summary_ar": "تسجيل عقد الزواج المدني المبرم في الخارج أو أمام كاتب العدل في السجلات اللبنانية.",
        "duration_ar": "2-4 أسابيع",
        "fee_ar": "رسوم دمغة 25,000 ل.ل.",
        "tags": ["زواج", "عقد", "مدني", "أحوال شخصية"],
        "source_tier": 2,
        "last_verified": "2024-01-01",
        "review_expiry": "2025-01-01",
        "status": "verified",
        "documents": [
            {"name_ar": "عقد الزواج الأصلي مع الترجمة", "required": True},
            {"name_ar": "بطاقة هوية الزوجين", "required": True},
            {"name_ar": "شهادة الميلاد للزوجين", "required": True},
            {"name_ar": "تصريح قيد نفوس", "required": True},
        ],
        "steps": [
            {"order": 1, "title_ar": "التحقق من عقد الزواج", "desc_ar": "التحقق من صحة عقد الزواج وتصديقه من السفارة اللبنانية إن كان خارجياً."},
            {"order": 2, "title_ar": "ترجمة الوثائق", "desc_ar": "ترجمة جميع الوثائق إلى اللغة العربية من مترجم معتمد."},
            {"order": 3, "title_ar": "تقديم الطلب لدائرة النفوس", "desc_ar": "تقديم الملف الكامل لدائرة النفوس في وزارة الداخلية."},
            {"order": 4, "title_ar": "استلام القيد", "desc_ar": "استلام قيد الزواج في سجل الأحوال الشخصية."},
        ],
    },
    {
        "id": "land-registration-lb",
        "slug": "land-registration-lb",
        "title_ar": "تسجيل عقار أو نقل ملكية",
        "title_en": "Land & Property Registration",
        "category": "real_estate",
        "country": "lebanon",
        "authority": "المديرية العامة للشؤون العقارية",
        "authority_en": "General Directorate of Land Registry",
        "authority_url": "https://www.cadastre.gov.lb",
        "summary_ar": "تسجيل عقار لأول مرة أو نقل ملكيته في السجل العقاري اللبناني.",
        "duration_ar": "3-6 أشهر",
        "fee_ar": "5% من قيمة العقار + رسوم إدارية",
        "tags": ["عقار", "ملكية", "قيد", "طابو"],
        "source_tier": 1,
        "last_verified": "2024-01-01",
        "review_expiry": "2025-01-01",
        "status": "verified",
        "documents": [
            {"name_ar": "صك الملكية أو عقد البيع", "required": True},
            {"name_ar": "بطاقة هوية المتعاقدين", "required": True},
            {"name_ar": "رسم تقريبي للعقار مصدّق", "required": True},
            {"name_ar": "براءة ذمة ضريبية", "required": True},
            {"name_ar": "مستخرج قيد للعقار", "required": True},
            {"name_ar": "تفويض كاتب عدل إن وجد وكيل", "required": False},
        ],
        "steps": [
            {"order": 1, "title_ar": "التحقق من صحة المستندات", "desc_ar": "التحقق من صحة عقد البيع وتصديقه من كاتب العدل."},
            {"order": 2, "title_ar": "سداد الضرائب والرسوم", "desc_ar": "سداد ضريبة انتقال الملكية 5% وكل الرسوم المستحقة."},
            {"order": 3, "title_ar": "تقديم الطلب لدائرة المساحة", "desc_ar": "تقديم الملف الكامل لمديرية الشؤون العقارية في المنطقة."},
            {"order": 4, "title_ar": "مراجعة الطلب والمعاينة", "desc_ar": "مراجعة الطلب من قبل الدائرة وإجراء معاينة ميدانية للعقار إن لزم."},
            {"order": 5, "title_ar": "قيد الملكية", "desc_ar": "قيد الملكية في السجل العقاري وإصدار صك الملكية الجديد."},
        ],
    },
    {
        "id": "work-permit-lb",
        "slug": "work-permit-lb",
        "title_ar": "استخراج إذن عمل للأجانب",
        "title_en": "Work Permit for Foreigners",
        "category": "education",
        "country": "lebanon",
        "authority": "وزارة العمل",
        "authority_en": "Ministry of Labor",
        "authority_url": "https://www.labor.gov.lb",
        "summary_ar": "الحصول على إذن عمل رسمي للعمال الأجانب في لبنان وتجديده سنوياً.",
        "duration_ar": "4-8 أسابيع",
        "fee_ar": "800,000 ل.ل. سنوياً",
        "tags": ["عمل", "أجنبي", "إذن", "تأشيرة عمل"],
        "source_tier": 2,
        "last_verified": "2024-01-01",
        "review_expiry": "2025-01-01",
        "status": "verified",
        "documents": [
            {"name_ar": "جواز سفر ساري المفعول", "required": True},
            {"name_ar": "عقد عمل مصدق من وزارة العمل", "required": True},
            {"name_ar": "شهادة طبية مصدقة", "required": True},
            {"name_ar": "صور شمسية", "required": True},
            {"name_ar": "شهادة الكفيل أو صاحب العمل", "required": True},
            {"name_ar": "شهادة الخلو من السوابق من بلد المنشأ", "required": False},
        ],
        "steps": [
            {"order": 1, "title_ar": "تصديق عقد العمل", "desc_ar": "تصديق عقد العمل من صاحب العمل وتوثيقه لدى كاتب العدل."},
            {"order": 2, "title_ar": "تقديم الطلب لوزارة العمل", "desc_ar": "تقديم كامل الملف لوزارة العمل مع رسم الطلب."},
            {"order": 3, "title_ar": "فحص الطلب", "desc_ar": "فحص الطلب من قبل اللجنة المختصة والتحقق من عدم وجود لبناني يشغل الوظيفة."},
            {"order": 4, "title_ar": "سداد الرسوم", "desc_ar": "سداد رسوم إذن العمل السنوية."},
            {"order": 5, "title_ar": "استلام إذن العمل", "desc_ar": "استلام إذن العمل ومنه التسجيل في الضمان الاجتماعي."},
        ],
    },
]

def _proc_index() -> dict:
    return {p["id"]: p for p in _PROCEDURE_CATALOG}

@app.get("/procedures")
async def list_procedures(
    country: Optional[str] = None,
    category: Optional[str] = None,
    q: Optional[str] = None,
    limit: int = 20,
):
    """List government procedures with optional filtering."""
    results = _PROCEDURE_CATALOG.copy()
    if country:
        results = [p for p in results if p.get("country", "").lower() == country.lower()]
    if category:
        results = [p for p in results if p.get("category", "").lower() == category.lower()]
    if q:
        q_lower = q.lower()
        results = [p for p in results if (
            q_lower in p.get("title_ar", "").lower() or
            q_lower in p.get("title_en", "").lower() or
            any(q_lower in t for t in p.get("tags", []))
        )]
    return {
        "total": len(results),
        "procedures": [
            {
                "id": p["id"],
                "slug": p["slug"],
                "title_ar": p["title_ar"],
                "title_en": p.get("title_en", ""),
                "category": p["category"],
                "country": p["country"],
                "authority": p["authority"],
                "summary_ar": p.get("summary_ar", ""),
                "status": p.get("status", "unverified"),
                "last_verified": p.get("last_verified", ""),
            }
            for p in results[:limit]
        ],
    }

@app.get("/procedures/categories/list")
async def procedure_categories():
    """Return distinct categories."""
    cats = list({p["category"] for p in _PROCEDURE_CATALOG})
    return {"categories": cats}

@app.get("/procedures/{procedure_id}")
async def get_procedure(procedure_id: str):
    """Get full procedure detail by ID or slug."""
    idx = _proc_index()
    # Try by ID, then by slug
    proc = idx.get(procedure_id)
    if not proc:
        proc = next((p for p in _PROCEDURE_CATALOG if p.get("slug") == procedure_id), None)
    if not proc:
        raise HTTPException(404, detail="الإجراء غير موجود")
    return proc

# ══════════════════════════════════════════════════════════════
#  MY PROCEDURE WORKSPACE — per-user transaction tracker
#  Storage: in-memory dict (ephemeral; PostgreSQL migration BLOCKED
#  until owner confirms backup + restore test).
# ══════════════════════════════════════════════════════════════


async def _require_auth(current_user: dict = Depends(get_current_user)) -> dict:
    """Raise 401 if the request carries no real JWT (i.e. guest access)."""
    if current_user.get("username") == "guest":
        raise HTTPException(status_code=401, detail="يلزم تسجيل الدخول أولاً")
    return current_user


class MyProcedureCreate(BaseModel):
    procedure_id: str
    title_ar: Optional[str] = None
    notes: Optional[str] = None

class MyProcedureUpdate(BaseModel):
    notes: Optional[str] = None
    completed_steps: Optional[list[int]] = None
    status: Optional[str] = None  # active | completed | cancelled

# In-memory workspace (ephemeral — survives restart only if persisted to Qdrant user payload)
_user_procedures: dict[str, list[dict]] = {}

def _get_user_procedures(username: str) -> list[dict]:
    return _user_procedures.get(username, [])

def _save_user_procedure(username: str, proc: dict) -> None:
    if username not in _user_procedures:
        _user_procedures[username] = []
    existing = next((p for p in _user_procedures[username] if p["id"] == proc["id"]), None)
    if existing:
        existing.update(proc)
    else:
        _user_procedures[username].append(proc)

@app.get("/my-procedures")
async def list_my_procedures(user: dict = Depends(_require_auth)):
    procs = _get_user_procedures(user["username"])
    return {"procedures": procs, "count": len(procs)}

@app.post("/my-procedures")
async def create_my_procedure(req: MyProcedureCreate, user: dict = Depends(_require_auth)):
    """Start tracking a procedure for the current user."""
    # Find the procedure template
    idx = _proc_index()
    template = idx.get(req.procedure_id)
    if not template:
        template = next((p for p in _PROCEDURE_CATALOG if p.get("slug") == req.procedure_id), None)

    proc_id = f"{user['username']}_{req.procedure_id}_{uuid.uuid4().hex[:8]}"

    # Build checklist from template steps
    if template:
        checklist = [
            {
                "step": step["order"],
                "title_ar": step.get("title_ar", ""),
                "desc_ar": step.get("desc_ar", ""),
                "done": False,
                "done_at": None,
            }
            for step in template.get("steps", [])
        ]
        doc_list = [
            {
                "name_ar": doc["name_ar"],
                "required": doc.get("required", True),
                "uploaded": False,
            }
            for doc in template.get("documents", [])
        ]
    else:
        checklist = []
        doc_list = []

    user_proc = {
        "id": proc_id,
        "procedure_id": req.procedure_id,
        "title_ar": req.title_ar or (template["title_ar"] if template else req.procedure_id),
        "status": "active",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "updated_at": datetime.now(timezone.utc).isoformat(),
        "notes": req.notes or "",
        "checklist": checklist,
        "documents": doc_list,
        "completion_pct": 0,
        "next_step": checklist[0]["title_ar"] if checklist else "",
    }
    _save_user_procedure(user["username"], user_proc)
    return user_proc

@app.get("/my-procedures/{proc_id}")
async def get_my_procedure(proc_id: str, user: dict = Depends(_require_auth)):
    procs = _get_user_procedures(user["username"])
    proc = next((p for p in procs if p["id"] == proc_id), None)
    if not proc:
        raise HTTPException(404, detail="الملف غير موجود")
    return proc

@app.put("/my-procedures/{proc_id}")
async def update_my_procedure(
    proc_id: str, req: MyProcedureUpdate, user: dict = Depends(_require_auth)
):
    procs = _get_user_procedures(user["username"])
    proc = next((p for p in procs if p["id"] == proc_id), None)
    if not proc:
        raise HTTPException(404, detail="الملف غير موجود")

    if req.notes is not None:
        proc["notes"] = req.notes
    if req.status is not None:
        proc["status"] = req.status
    if req.completed_steps is not None:
        for step in proc.get("checklist", []):
            step["done"] = step["step"] in req.completed_steps
            if step["done"] and not step.get("done_at"):
                step["done_at"] = datetime.now(timezone.utc).isoformat()
        total = len(proc["checklist"])
        done = sum(1 for s in proc["checklist"] if s["done"])
        proc["completion_pct"] = round((done / total) * 100) if total else 0
        next_steps = [s for s in proc["checklist"] if not s["done"]]
        proc["next_step"] = next_steps[0]["title_ar"] if next_steps else "مكتمل"

    proc["updated_at"] = datetime.now(timezone.utc).isoformat()
    _save_user_procedure(user["username"], proc)
    return proc

@app.delete("/my-procedures/{proc_id}")
async def delete_my_procedure(proc_id: str, user: dict = Depends(_require_auth)):
    if user["username"] in _user_procedures:
        _user_procedures[user["username"]] = [
            p for p in _user_procedures[user["username"]] if p["id"] != proc_id
        ]
    return {"message": "تم حذف الملف"}


# ─────────────────────────────────────────────────────────────────────────────
# Admin Content Governance — Draft → Review → Approved → Published → Expired
# ─────────────────────────────────────────────────────────────────────────────

# In-memory content items store (interim; will move to PostgreSQL)
_content_items: dict[str, dict] = {}
_audit_log: list[dict] = []

_VALID_CONTENT_STATUSES = {"draft", "review", "approved", "published", "expired"}

_CONTENT_TRANSITIONS: dict[str, set] = {
    "draft":     {"review"},
    "review":    {"approved", "draft"},
    "approved":  {"published", "review"},
    "published": {"expired"},
    "expired":   set(),
}


def _require_admin(current_user: dict = Depends(get_current_user)) -> dict:
    """Raise 403 unless the current user has admin role or plan."""
    if current_user.get("plan") != "admin" and current_user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="مخصص للمسؤولين فقط")
    return current_user


def _append_audit(
    *,
    action: str,
    item_id: str,
    actor: str,
    before: str | None = None,
    after: str | None = None,
    note: str = "",
) -> None:
    _audit_log.append(
        {
            "ts": datetime.utcnow().isoformat(),
            "action": action,
            "item_id": item_id,
            "actor": actor,
            "before": before,
            "after": after,
            "note": note,
        }
    )


class ContentItemCreate(BaseModel):
    title_ar: str
    body_ar: str
    content_type: str = "procedure_update"
    ref_procedure_id: str | None = None
    note: str = ""


class ContentItemPatch(BaseModel):
    title_ar: str | None = None
    body_ar: str | None = None
    note: str = ""


class ContentStatusTransition(BaseModel):
    target_status: str
    note: str = ""


# ── List content items (admin only) ──────────────────────────────────────────
@app.get("/admin/content", tags=["admin"])
async def admin_list_content(
    status: str | None = None,
    admin: dict = Depends(_require_admin),
):
    items = list(_content_items.values())
    if status:
        items = [i for i in items if i["status"] == status]
    items.sort(key=lambda x: x["updated_at"], reverse=True)
    return {"items": items, "count": len(items)}


# ── Create draft content item ─────────────────────────────────────────────────
@app.post("/admin/content", tags=["admin"], status_code=201)
async def admin_create_content(
    body: ContentItemCreate,
    admin: dict = Depends(_require_admin),
):
    item_id = secrets.token_hex(8)
    now = datetime.utcnow().isoformat()
    item = {
        "id": item_id,
        "title_ar": body.title_ar,
        "body_ar": body.body_ar,
        "content_type": body.content_type,
        "ref_procedure_id": body.ref_procedure_id,
        "status": "draft",
        "created_by": admin["username"],
        "created_at": now,
        "updated_at": now,
    }
    _content_items[item_id] = item
    _append_audit(
        action="CREATE",
        item_id=item_id,
        actor=admin["username"],
        after="draft",
        note=body.note,
    )
    return item


# ── Get single content item ───────────────────────────────────────────────────
@app.get("/admin/content/{item_id}", tags=["admin"])
async def admin_get_content(item_id: str, admin: dict = Depends(_require_admin)):
    if item_id not in _content_items:
        raise HTTPException(status_code=404, detail="العنصر غير موجود")
    return _content_items[item_id]


# ── Update draft content item ─────────────────────────────────────────────────
@app.patch("/admin/content/{item_id}", tags=["admin"])
async def admin_update_content(
    item_id: str,
    body: ContentItemPatch,
    admin: dict = Depends(_require_admin),
):
    if item_id not in _content_items:
        raise HTTPException(status_code=404, detail="العنصر غير موجود")
    item = _content_items[item_id]
    if item["status"] not in {"draft", "review"}:
        raise HTTPException(
            status_code=409,
            detail=f"لا يمكن تعديل المحتوى في حالة {item['status']}",
        )
    if body.title_ar is not None:
        item["title_ar"] = body.title_ar
    if body.body_ar is not None:
        item["body_ar"] = body.body_ar
    item["updated_at"] = datetime.utcnow().isoformat()
    _append_audit(
        action="EDIT",
        item_id=item_id,
        actor=admin["username"],
        before=item["status"],
        after=item["status"],
        note=body.note,
    )
    return item


# ── Transition content status ─────────────────────────────────────────────────
@app.post("/admin/content/{item_id}/transition", tags=["admin"])
async def admin_transition_content(
    item_id: str,
    body: ContentStatusTransition,
    admin: dict = Depends(_require_admin),
):
    if item_id not in _content_items:
        raise HTTPException(status_code=404, detail="العنصر غير موجود")
    item = _content_items[item_id]
    current = item["status"]
    target = body.target_status.lower()

    if target not in _VALID_CONTENT_STATUSES:
        raise HTTPException(status_code=400, detail=f"الحالة '{target}' غير صالحة")
    if target not in _CONTENT_TRANSITIONS.get(current, set()):
        raise HTTPException(
            status_code=409,
            detail=f"الانتقال من '{current}' إلى '{target}' غير مسموح",
        )

    item["status"] = target
    item["updated_at"] = datetime.utcnow().isoformat()
    if target == "published":
        item["published_at"] = item["updated_at"]

    _append_audit(
        action="TRANSITION",
        item_id=item_id,
        actor=admin["username"],
        before=current,
        after=target,
        note=body.note,
    )
    return {"id": item_id, "status": item["status"], "updated_at": item["updated_at"]}


# ── Audit log ─────────────────────────────────────────────────────────────────
@app.get("/admin/audit-log", tags=["admin"])
async def admin_audit_log(
    limit: int = 50,
    offset: int = 0,
    admin: dict = Depends(_require_admin),
):
    log = list(reversed(_audit_log))
    total = len(log)
    return {
        "entries": log[offset : offset + limit],
        "total": total,
        "limit": limit,
        "offset": offset,
    }

# ═══════════════════════════════════════════════════════════════
#  STARTUP
# ═══════════════════════════════════════════════════════════════

@app.on_event("startup")
async def startup():
    admin_username = os.environ.get("ADMIN_USERNAME")
    admin_password = os.environ.get("ADMIN_PASSWORD")
    if admin_username and admin_password:
        try:
            existing = db_get_user(admin_username.lower())
            if not existing:
                now = datetime.now(timezone.utc)
                admin_user = {
                    "username":         admin_username.lower(),
                    "email":            os.environ.get("ADMIN_EMAIL", f"{admin_username.lower()}@dalilak.ai"),
                    "password_hash":    hash_pw(admin_password),
                    "full_name":        "مشرف النظام",
                    "phone":            "",
                    "plan":             "admin",
                    "role":             "admin",
                    "active":           True,
                    "trial_expires_at": None,
                    "paid_until":       None,
                    "created_at":       now.isoformat(),
                    "last_login":       None,
                }
                db_save_user(admin_user)
                _log.info("[startup] Admin user '%s' created", admin_username)
            else:
                _log.info("[startup] Admin user '%s' already exists", admin_username)
        except Exception as e:
            _log.error("[startup] Failed to create admin user: %s", e)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
